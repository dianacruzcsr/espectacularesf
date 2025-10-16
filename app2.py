import streamlit as st
import pandas as pd
from pptx import Presentation
from geopy.distance import geodesic
from datetime import datetime
import folium
from folium.plugins import MarkerCluster
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from copy import deepcopy
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import re
import io

# ================================
# Funciones de lógica de negocio
# ================================

def generar_folio(lugares_multiples=None):
    """Genera un folio único con formato NEGOCIO-FECHA"""
    try:
        ahora = datetime.now()
        
        # Verificar si session_state existe y tiene nombre_negocio
        if (hasattr(st, 'session_state') and 
            hasattr(st.session_state, 'nombre_negocio') and 
            st.session_state.nombre_negocio and 
            st.session_state.nombre_negocio.strip()):
            
            nombre_negocio = st.session_state.nombre_negocio.strip()
            nombre_limpio = re.sub(r'[^a-zA-Z0-9]', '', nombre_negocio)[:15] if nombre_negocio else "NEGOCIO"
        else:
            # Fallback al primer lugar si no hay nombre de negocio
            if lugares_multiples and len(lugares_multiples) > 0:
                primer_lugar = lugares_multiples[0].get("nombre", "LUGAR")
                nombre_limpio = re.sub(r'[^a-zA-Z0-9]', '', primer_lugar)[:15] if primer_lugar else "LUGAR"
            else:
                nombre_limpio = "NEGOCIO"
        
        folio = f"{nombre_limpio}-{ahora.year}{ahora.month:02d}{ahora.day:02d}-{ahora.hour:02d}{ahora.minute:02d}"
        return folio
    except Exception as e:
        # Fallback seguro
        ahora = datetime.now()
        return f"NEGOCIO-{ahora.year}{ahora.month:02d}{ahora.day:02d}-{ahora.hour:02d}{ahora.minute:02d}"

def incrementar_folio(lugares_multiples=None):
    """Incrementa el contador de folios solo cuando se realiza una descarga"""
    if 'contador_descargas' not in st.session_state:
        st.session_state.contador_descargas = 1
    else:
        st.session_state.contador_descargas += 1
    
    # Verificar si podemos generar un folio con session_state
    try:
        st.session_state.folio_actual = generar_folio(lugares_multiples)
    except AttributeError:
        # Si hay error, generar folio básico
        ahora = datetime.now()
        st.session_state.folio_actual = f"NEGOCIO-{ahora.year}{ahora.month:02d}{ahora.day:02d}-{ahora.hour:02d}{ahora.minute:02d}"

def analizar_formato_coordenada(coord_str):
    if pd.isna(coord_str) or coord_str == "" or str(coord_str).strip() == "":
        return "vacía", None, None, None
    coord_str = str(coord_str).strip()
    patron_grados_dir = r'(\d+)°\s*(\d+)\'\s*(\d+\.?\d*)\"\s*([NSWE])'
    coincidencia = re.search(patron_grados_dir, coord_str, re.IGNORECASE)
    if coincidencia:
        grados = float(coincidencia.group(1))
        minutos = float(coincidencia.group(2))
        segundos = float(coincidencia.group(3))
        direccion = coincidencia.group(4).upper()
        decimal = grados + minutos/60 + segundos/3600
        if direccion in ['S', 'W']: decimal = -decimal
        return "grados_dir", decimal, direccion, None
    patron_grados_sig = r'([+-]?\d+\.\d+)'
    coincidencia = re.search(patron_grados_sig, coord_str)
    if coincidencia:
        decimal = float(coincidencia.group(1))
        return "grados_sig", decimal, None, None
    patron_dms = r'([+-]?\d+)\s+(\d+)\s+(\d+\.?\d*)'
    coincidencia = re.search(patron_dms, coord_str)
    if coincidencia:
        grados = float(coincidencia.group(1))
        minutos = float(coincidencia.group(2))
        segundos = float(coincidencia.group(3))
        decimal = abs(grados) + minutos/60 + segundos/3600
        if grados < 0: decimal = -decimal
        return "dms", decimal, None, None
    patron_dm = r'([+-]?\d+)°\s*(\d+\.\d+)'
    coincidencia = re.search(patron_dm, coord_str)
    if coincidencia:
        grados = float(coincidencia.group(1))
        minutos = float(coincidencia.group(2))
        decimal = abs(grados) + minutos/60
        if grados < 0: decimal = -decimal
        return "dm", decimal, None, None
    if "," in coord_str and "." not in coord_str:
        try:
            coord_europeo = coord_str.replace(",", ".")
            decimal = float(coord_europeo)
            return "decimal_eu", decimal, None, None
        except ValueError:
            pass
    if coord_str.count(",") >= 2:
        try:
            coord_sin_comas = coord_str.replace(",", "")
            decimal = float(coord_sin_comas)
            return "decimal_miles", decimal, None, None
        except ValueError:
            pass
    try:
        decimal = float(coord_str)
        return "decimal", decimal, None, None
    except ValueError:
        pass
    return "desconocido", None, None, None

def ajustar_valor_utm(valor, es_latitud=True):
    if valor is None: return None
    valor_abs = abs(valor)
    if es_latitud and (-90 <= valor <= 90): return valor
    elif not es_latitud and (-180 <= valor <= 180): return valor
    if valor_abs > 180:
        if valor_abs >= 1000000000: divisor = 10000000
        elif valor_abs >= 100000000: divisor = 1000000
        elif valor_abs >= 10000000: divisor = 100000
        elif valor_abs >= 1000000: divisor = 10000
        elif valor_abs >= 100000: divisor = 1000
        else: divisor = 1000
        valor_ajustado = valor / divisor
        if es_latitud and (-90 <= valor_ajustado <= 90): return valor_ajustado
        elif not es_latitud and (-180 <= valor_ajustado <= 180): return valor_ajustado
        else: return valor / (divisor * 10)
    return valor

def estandarizar_coordenada_universal(coord_str, formato, valor_raw=None, direccion=None, es_latitud=True):
    if pd.isna(coord_str) or coord_str == "": return None
    coord_str = str(coord_str).strip()
    try:
        if formato in ["grados_dir", "grados_sig", "dms", "dm"] and valor_raw is not None:
            valor = valor_raw
        elif formato == "decimal_eu":
            valor = float(coord_str.replace(",", "."))
        elif formato == "decimal_miles":
            valor = float(coord_str.replace(",", ""))
        elif formato == "decimal":
            valor = float(coord_str)
        elif formato == "desconocido":
            numeros = re.findall(r'[+-]?\d+\.?\d*', coord_str)
            valor = float(numeros[0]) if numeros else None
        else:
            valor = None
        if valor is None: return None
        return ajustar_valor_utm(valor, es_latitud)
    except (ValueError, TypeError):
        return None

def detectar_inversion_universal(lat_str, lon_str):
    formato_lat, valor_lat, dir_lat, _ = analizar_formato_coordenada(lat_str)
    formato_lon, valor_lon, dir_lon, _ = analizar_formato_coordenada(lon_str)
    if valor_lat is None or valor_lon is None: return False, formato_lat, formato_lon, ["no_analizable"]
    valor_abs_lat, valor_abs_lon = abs(valor_lat), abs(valor_lon)
    es_patron_utm_invertido = ((valor_abs_lat > 1000000 or valor_abs_lon > 1000000) and
                              (valor_abs_lat < 1000000000 and valor_abs_lon < 1000000000) and
                              (valor_lat < 0 and valor_lon > 0))
    if es_patron_utm_invertido: return True, formato_lat, formato_lon, ["patron_utm_invertido"]
    rango_lat_absoluto, rango_lon_absoluto = (-90, 90), (-180, 180)
    lat_en_rango_valido = rango_lat_absoluto[0] <= valor_lat <= rango_lat_absoluto[1]
    lon_en_rango_valido = rango_lon_absoluto[0] <= valor_lon <= rango_lon_absoluto[1]
    direcciones_invertidas = (dir_lat in ['E', 'W'] and dir_lon in ['N', 'S']) if dir_lat and dir_lon else False
    lat_podria_ser_lon = (rango_lon_absoluto[0] <= valor_lat <= rango_lon_absoluto[1] and not lat_en_rango_valido)
    lon_podria_ser_lat = (rango_lat_absoluto[0] <= valor_lon <= rango_lat_absoluto[1] and not lon_en_rango_valido)
    diferencia_magnitud = abs(abs(valor_lat) - abs(valor_lon))
    lat_mayor_que_lon = abs(valor_lat) > abs(valor_lon) + 10
    signos_atipicos = (valor_lat < 0 and valor_lon > 0)
    criterios = []
    if direcciones_invertidas: criterios.append("direcciones_cardinales")
    if not lat_en_rango_valido and not lon_en_rango_valido and lat_podria_ser_lon and lon_podria_ser_lat: criterios.append("valores_fuera_de_rango")
    if lat_mayor_que_lon and diferencia_magnitud > 20: criterios.append("diferencia_magnitud")
    if signos_atipicos: criterios.append("signos_atipicos")
    probable_inversion = len(criterios) >= 2
    return probable_inversion, formato_lat, formato_lon, criterios

def duplicar_slide(prs, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)
    for shp in slide.shapes:
        el = deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    for rel in slide.part.rels.values():
        try:
            if rel.reltype == RT.IMAGE:
                image_part = rel.target_part
                new_slide.part.relate_to(image_part, RT.IMAGE)
            elif rel.reltype == RT.HYPERLINK:
                new_slide.part.relate_to(rel.target_ref, RT.HYPERLINK)
        except Exception as e:
            st.error(f"⚠️ No se pudo copiar relación: {e}")
    return new_slide

def reemplazar_texto_slide(slide, fila, df_filtrado, nombre_negocio=""):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if nombre_negocio and "{{NOMBRE_NEGOCIO}}" in run.text:
                    run.text = run.text.replace("{{NOMBRE_NEGOCIO}}", nombre_negocio)
                
                for campo in df_filtrado.columns:
                    marcador = f"{{{{{campo.strip()}}}}}"
                    
                    if marcador in run.text:
                        valor = fila.get(campo, "")
                        
                        if campo.strip().upper() == "LATITUD":
                            try:
                                run.text = run.text.replace(marcador, str(valor))
                                run.hyperlink.address = fila.get("STREET_VIEW", "#")
                            except Exception as e:
                                st.warning(f"⚠️ No se pudo crear el hipervínculo para {campo}: {e}")
                        else:
                            run.text = run.text.replace(marcador, str(valor))

def procesar_busqueda_individual(df_uploaded, lat_negocio, lon_negocio, radio_km, presupuesto_min, presupuesto_max, tipos_seleccionados, nombre_lugar=""):
    with st.spinner(f"🔄 Analizando coordenadas para {nombre_lugar if nombre_lugar else 'el lugar'}..."):
        df_copy = df_uploaded.copy()
        
        if tipos_seleccionados and "TIPO" in df_copy.columns:
            df_copy = df_copy[df_copy["TIPO"].isin(tipos_seleccionados)]
        
        latitudes_decimal = []
        longitudes_decimal = []
        
        for _, row in df_copy.iterrows():
            lat_original, lon_original = row["LATITUD"], row["LONGITUD"]
            formato_lat, valor_lat, dir_lat, _ = analizar_formato_coordenada(lat_original)
            formato_lon, valor_lon, dir_lon, _ = analizar_formato_coordenada(lon_original)
            invertidas, f_lat_det, f_lon_det, criterios = detectar_inversion_universal(lat_original, lon_original)
            
            if invertidas:
                lat_corregida = estandarizar_coordenada_universal(lon_original, f_lon_det, valor_lon, dir_lon, True)
                lon_corregida = estandarizar_coordenada_universal(lat_original, f_lat_det, valor_lat, dir_lat, False)
            else:
                lat_corregida = estandarizar_coordenada_universal(lat_original, formato_lat, valor_lat, dir_lat, True)
                lon_corregida = estandarizar_coordenada_universal(lon_original, formato_lon, valor_lon, dir_lon, False)
            
            latitudes_decimal.append(lat_corregida)
            longitudes_decimal.append(lon_corregida)
    
    # CORRECCIÓN: Usar listas directamente en lugar de zip
    df_copy["LATITUD_DECIMAL"] = latitudes_decimal
    df_copy["LONGITUD_DECIMAL"] = longitudes_decimal
    
    resultados = []
    for _, row in df_copy.iterrows():
        try:
            lat_raw, lon_raw = row["LATITUD_DECIMAL"], row["LONGITUD_DECIMAL"]
            
            # VERIFICACIÓN MÁS ROBUSTA
            if (pd.isna(lat_raw) or pd.isna(lon_raw) or 
                not isinstance(lat_raw, (int, float)) or 
                not isinstance(lon_raw, (int, float)) or
                not (-90 <= lat_raw <= 90) or 
                not (-180 <= lon_raw <= 180)): 
                continue
                
            distancia = geodesic((lat_negocio, lon_negocio), (lat_raw, lon_raw)).km
            if distancia < radio_km:
                tarifa_val = pd.to_numeric(row["TARIFA"], errors="coerce")
                if (presupuesto_min is not None and tarifa_val < presupuesto_min) or \
                   (presupuesto_max is not None and tarifa_val > presupuesto_max): 
                    continue
                
                maps_url = f"https://www.google.com/maps/place/{lat_raw},{lon_raw}"
                street_view_url = f"https://www.google.com/maps/@?api=1&map_action=pano&viewpoint={lat_raw},{lon_raw}"

                resultados.append({
                    "CIUDAD": row.get("CIUDAD"), "CLAVE": row.get("CLAVE"), "DIRECCION": row.get("DIRECCION"),
                    "VISTA": row.get("VISTA"), "TIPO": row.get("TIPO"), "BASE": row.get("BASE"),
                    "ALTURA": row.get("ALTURA"), "AREA": row.get("AREA"), "LATITUD": lat_raw,
                    "LONGITUD": lon_raw, "DISTANCIA_KM": round(distancia, 2),
                    "TARIFA_PUBLICO": f"${tarifa_val:,.2f}", "IMPRESION": row.get("IMPRESION"),
                    "INSTALACION": row.get("INSTALACION"), "COSTO": row.get("IMPRESION+INSTALACION"),
                    "MAPS_": maps_url,
                    "STREET_VIEW": street_view_url,
                    "PROVEEDOR": row.get("PROVEEDOR"), "TELEFONO_PROVEEDOR": row.get("TELÉFONO PROVEEDOR"),
                    "LUGAR_BUSQUEDA": nombre_lugar if nombre_lugar else "Principal",
                    "LAT_NEGOCIO": lat_negocio,
                    "LON_NEGOCIO": lon_negocio
                })
        except Exception as e:
            st.warning(f"Error al procesar fila: {e}")
    
    return pd.DataFrame(resultados)

# ================================
# ESTRUCTURA DE LA APP STREAMLIT
# ================================

st.set_page_config(layout="wide")
st.title("🤖 Analizador de Espectaculares")
st.subheader("¡Hola! Soy tu asistente virtual para encontrar los mejores espectaculares publicitarios para tu negocio.")
st.write("---")

# Inicializar st.session_state
if 'df_filtrado' not in st.session_state:
    st.session_state.df_filtrado = pd.DataFrame()
if 'uploaded_df' not in st.session_state:
    st.session_state.uploaded_df = None
if 'negocio_lat' not in st.session_state:
    st.session_state.negocio_lat = 19.4326
if 'negocio_lon' not in st.session_state:
    st.session_state.negocio_lon = -99.1332
if 'presupuesto_min' not in st.session_state:
    st.session_state.presupuesto_min = 0.0
if 'presupuesto_max' not in st.session_state:
    st.session_state.presupuesto_max = 100000.0
if 'radio_km' not in st.session_state:
    st.session_state.radio_km = 5.0
if 'contador_descargas' not in st.session_state:
    st.session_state.contador_descargas = 1
if 'folio_actual' not in st.session_state:
    st.session_state.folio_actual = "NEGOCIO-" + datetime.now().strftime("%Y%m%d-%H%M")
if 'busqueda_realizada' not in st.session_state:
    st.session_state.busqueda_realizada = False
if 'espectaculares_seleccionados' not in st.session_state:
    st.session_state.espectaculares_seleccionados = []
if 'nombre_negocio' not in st.session_state:
    st.session_state.nombre_negocio = ""
if 'tipos_espectaculares' not in st.session_state:
    st.session_state.tipos_espectaculares = []
if 'tipos_seleccionados' not in st.session_state:
    st.session_state.tipos_seleccionados = []
if 'consultas_previas' not in st.session_state:
    st.session_state.consultas_previas = []
if 'lugares_multiples' not in st.session_state:
    st.session_state.lugares_multiples = [{"nombre": "Principal", "lat": 19.4326, "lon": -99.1332}]
if 'busqueda_combinada' not in st.session_state:
    st.session_state.busqueda_combinada = pd.DataFrame()
if 'indices_seleccionados' not in st.session_state:
    st.session_state.indices_seleccionados = []
if 'df_por_lugar' not in st.session_state:
    st.session_state.df_por_lugar = {}
if 'selecciones_por_lugar' not in st.session_state:
    st.session_state.selecciones_por_lugar = {}
if 'multiselect_actualizado' not in st.session_state:
    st.session_state.multiselect_actualizado = False

# 1. UPLOAD CSV
uploaded_file = st.file_uploader("📂 **Paso 1: Sube tu archivo CSV de inventario**", type="csv")
if uploaded_file:
    st.session_state.uploaded_df = pd.read_csv(uploaded_file, sep=",")
    st.session_state.uploaded_df.columns = st.session_state.uploaded_df.columns.str.strip()
    st.success(f"✅ CSV cargado con **{len(st.session_state.uploaded_df)}** registros.")
    if "TARIFA PUBLICO" not in st.session_state.uploaded_df.columns:
        st.error("❌ No se encuentra la columna **'TARIFA PUBLICO'** en el CSV.")
        st.session_state.uploaded_df = None
        st.stop()
    else:
        st.session_state.uploaded_df["TARIFA"] = (st.session_state.uploaded_df["TARIFA PUBLICO"].astype(str).str.replace(r"[^\d.]", "", regex=True).replace("", "0").astype(float))
    
    if "TIPO" in st.session_state.uploaded_df.columns:
        tipos_unicos = sorted(st.session_state.uploaded_df["TIPO"].dropna().unique().tolist())
        st.session_state.tipos_espectaculares = tipos_unicos
        if not st.session_state.tipos_seleccionados:
            st.session_state.tipos_seleccionados = tipos_unicos

# 2. INPUTS PARA LA BÚSQUEDA
st.write("---")
st.header("🎯 **Paso 2: Define tus criterios de búsqueda**")

st.session_state.nombre_negocio = st.text_input(
    "🏢 **Nombre del negocio:**", 
    value=st.session_state.nombre_negocio,
    placeholder="Ingresa el nombre de tu negocio..."
)

st.subheader("📍 Configuración de Lugares a Buscar")

for i, lugar in enumerate(st.session_state.lugares_multiples):
    col1, col2, col3, col4 = st.columns([2, 2, 2, 1])
    
    with col1:
        nombre_lugar = st.text_input(
            f"Nombre lugar {i+1}",
            value=lugar["nombre"],
            key=f"nombre_lugar_{i}"
        )
        st.session_state.lugares_multiples[i]["nombre"] = nombre_lugar
    
    with col2:
        lat_lugar = st.text_input(
            f"Latitud {i+1}",
            value=str(lugar["lat"]),
            key=f"lat_lugar_{i}"
        )
        try:
            st.session_state.lugares_multiples[i]["lat"] = float(lat_lugar) if lat_lugar.strip() != "" else 19.4326
        except ValueError:
            st.session_state.lugares_multiples[i]["lat"] = 19.4326
    
    with col3:
        lon_lugar = st.text_input(
            f"Longitud {i+1}",
            value=str(lugar["lon"]),
            key=f"lon_lugar_{i}"
        )
        try:
            st.session_state.lugares_multiples[i]["lon"] = float(lon_lugar) if lon_lugar.strip() != "" else -99.1332
        except ValueError:
            st.session_state.lugares_multiples[i]["lon"] = -99.1332
    
    with col4:
        if i > 0:
            if st.button("❌", key=f"eliminar_{i}"):
                st.session_state.lugares_multiples.pop(i)
                st.rerun()

if st.button("➕ Agregar otro lugar"):
    st.session_state.lugares_multiples.append({
        "nombre": f"Lugar {len(st.session_state.lugares_multiples) + 1}",
        "lat": 19.4326,
        "lon": -99.1332
    })
    st.rerun()

col1, col2, col3 = st.columns(3)

with col1:
    st.session_state.presupuesto_min = st.number_input("💰 Presupuesto mínimo:", value=st.session_state.presupuesto_min, format="%.2f", key='presupuesto_min_input')

with col2:
    st.session_state.presupuesto_max = st.number_input("💰 Presupuesto máximo:", value=st.session_state.presupuesto_max, format="%.2f", key='presupuesto_max_input')

with col3:
    st.session_state.radio_km = st.slider("📏 Radio de búsqueda (km):", min_value=0.5, max_value=50.0, value=st.session_state.radio_km, step=0.5, key='radio_km_input')

# Filtro por tipos de espectaculares
st.write("---")
st.subheader("🎪 **Selección de Tipos de Espectaculares**")

if st.session_state.tipos_espectaculares:
    col_btns1, col_btns2, col_btns3 = st.columns(3)
    
    with col_btns1:
        if st.button("✅ Seleccionar Todos", key="select_all"):
            st.session_state.tipos_seleccionados = st.session_state.tipos_espectaculares.copy()
            st.rerun()
    
    with col_btns2:
        if st.button("❌ Deseleccionar Todos", key="deselect_all"):
            st.session_state.tipos_seleccionados = []
            st.rerun()
    
    with col_btns3:
        if st.button("🔄 Invertir Selección", key="invert_selection"):
            if st.session_state.tipos_seleccionados:
                st.session_state.tipos_seleccionados = [tipo for tipo in st.session_state.tipos_espectaculares if tipo not in st.session_state.tipos_seleccionados]
            else:
                st.session_state.tipos_seleccionados = st.session_state.tipos_espectaculares.copy()
            st.rerun()
    
    st.session_state.tipos_seleccionados = st.multiselect(
        "**Selecciona uno o más tipos de espectaculares:**",
        options=st.session_state.tipos_espectaculares,
        default=st.session_state.tipos_seleccionados,
        placeholder="Elige los tipos de espectaculares que te interesan...",
        key='multiselect_tipos'
    )
    
    if st.session_state.tipos_seleccionados:
        if len(st.session_state.tipos_seleccionados) <= 5:
            st.success(f"✅ **{len(st.session_state.tipos_seleccionados)}** tipos seleccionados: {', '.join(st.session_state.tipos_seleccionados)}")
        else:
            st.success(f"✅ **{len(st.session_state.tipos_seleccionados)}** tipos seleccionados")
    else:
        st.info("ℹ️ No hay tipos seleccionados. Se mostrarán todos los espectaculares.")

# 3. FILTRADO Y GENERACIÓN DE RESULTADOS
if st.button("🚀 **Iniciar Búsqueda Múltiple**") and st.session_state.uploaded_df is not None:
    if len(st.session_state.lugares_multiples) == 0:
        st.warning("⚠️ Por favor, agrega al menos un lugar para buscar.")
    else:
        todos_resultados = []
        resultados_por_lugar = {}
        
        for lugar in st.session_state.lugares_multiples:
            with st.spinner(f"🔍 Buscando espectaculares cerca de: {lugar['nombre']}..."):
                df_resultado = procesar_busqueda_individual(
                    st.session_state.uploaded_df,
                    lugar["lat"],
                    lugar["lon"],
                    st.session_state.radio_km,
                    st.session_state.presupuesto_min,
                    st.session_state.presupuesto_max,
                    st.session_state.tipos_seleccionados,
                    lugar["nombre"]
                )
            
            if not df_resultado.empty:
                todos_resultados.append(df_resultado)
                resultados_por_lugar[lugar["nombre"]] = df_resultado  # ← AGREGAR ESTO
            else:
                st.warning(f"⚠️ **{lugar['nombre']}**: No se encontraron espectaculares")
        
        if todos_resultados:
            st.session_state.busqueda_combinada = pd.concat(todos_resultados, ignore_index=True)
            st.session_state.busqueda_combinada = st.session_state.busqueda_combinada.drop_duplicates(subset=['CLAVE'])
            
            st.session_state.df_filtrado = st.session_state.busqueda_combinada
            st.session_state.busqueda_realizada = True
            st.session_state.df_por_lugar = resultados_por_lugar  # ← GUARDAR RESULTADOS POR LUGAR
            
            st.session_state.folio_actual = generar_folio()  

            # REINICIAR LAS SELECCIONES AL HACER NUEVA BÚSQUEDA
            st.session_state.espectaculares_seleccionados = []
            st.session_state.indices_seleccionados = []
            st.session_state.selecciones_por_lugar = {}
            st.session_state.multiselect_actualizado = True
            
            consulta_actual = {
                "nombre_negocio": st.session_state.nombre_negocio,
                "fecha": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "resultados": len(st.session_state.busqueda_combinada),
                "df_filtrado": st.session_state.busqueda_combinada.copy(),
                "tipo": "múltiple",
                "lugares": [lugar["nombre"] for lugar in st.session_state.lugares_multiples],
                "tipos_seleccionados": st.session_state.tipos_seleccionados.copy() if st.session_state.tipos_seleccionados else []
            }
            st.session_state.consultas_previas.append(consulta_actual)
            
            st.success(f"🎉 **Búsqueda múltiple completada!** Se encontraron **{len(st.session_state.busqueda_combinada)}** espectaculares únicos cerca de {len(st.session_state.lugares_multiples)} lugares.")
        else:
            st.warning("⚠️ No se encontraron espectaculares en ninguno de los lugares especificados.")
            st.session_state.df_filtrado = pd.DataFrame()
            st.session_state.busqueda_realizada = False

# 4. VISUALIZACIÓN Y DESCARGAS
if not st.session_state.df_filtrado.empty and st.session_state.busqueda_realizada:
    df_filtrado = st.session_state.df_filtrado
    st.write("---")
    st.header("🔍 **Resultados de la Búsqueda Múltiple**")
    
    st.info(f"🏢 **Negocio:** {st.session_state.nombre_negocio if st.session_state.nombre_negocio else 'No especificado'} | 📋 **Folio:** `{st.session_state.folio_actual}`")
    st.info(f"📍 **Lugares buscados:** {', '.join([lugar['nombre'] for lugar in st.session_state.lugares_multiples])}")
    
    # DIAGNÓSTICO DE DATOS
    st.write("---")
    st.subheader("🔍 Diagnóstico de Datos para el Mapa")
    
    # Contar coordenadas válidas
    coordenadas_validas = 0
    coordenadas_invalidas = []
    
    for i, r in df_filtrado.iterrows():
        lat = r["LATITUD"]
        lon = r["LONGITUD"]
        
        if (pd.isna(lat) or pd.isna(lon) or 
            not isinstance(lat, (int, float)) or 
            not isinstance(lon, (int, float)) or
            not (-90 <= lat <= 90) or 
            not (-180 <= lon <= 180)):
            coordenadas_invalidas.append((i, r['CLAVE'], lat, lon))
        else:
            coordenadas_validas += 1
    
    col_diag1, col_diag2 = st.columns(2)
    
    with col_diag1:
        st.success(f"✅ **Coordenadas válidas:** {coordenadas_validas}")
    
    with col_diag2:
        if coordenadas_invalidas:
            st.error(f"❌ **Coordenadas inválidas:** {len(coordenadas_invalidas)}")
    
    if coordenadas_invalidas:
        with st.expander("📋 Ver coordenadas problemáticas"):
            for idx, clave, lat, lon in coordenadas_invalidas[:10]:  # Mostrar solo las primeras 10
                st.write(f"**Fila {idx}:** {clave} - Lat: `{lat}`, Lon: `{lon}`")
            if len(coordenadas_invalidas) > 10:
                st.info(f"... y {len(coordenadas_invalidas) - 10} más")
    
    st.subheader("🗺️ Mapa de Espectaculares (Todos los Lugares)")
    
    # Calcular centro del mapa de manera más robusta
    latitudes_validas = df_filtrado["LATITUD"].dropna()
    longitudes_validas = df_filtrado["LONGITUD"].dropna()
    
    if len(latitudes_validas) > 0 and len(longitudes_validas) > 0:
        centro_lat = latitudes_validas.mean()
        centro_lon = longitudes_validas.mean()
    else:
        centro_lat = 19.4326
        centro_lon = -99.1332
    
    mapa = folium.Map(location=[centro_lat, centro_lon], zoom_start=12)
    
    colores_lugares = ["red", "blue", "green", "purple", "orange", "darkred", "lightred", "beige", "darkblue", "darkgreen"]
    
    # Contadores para debugging
    marcadores_agregados = 0
    marcadores_fallados = 0
    
    # Primero agregar los lugares de búsqueda
    for i, lugar in enumerate(st.session_state.lugares_multiples):
        color = colores_lugares[i % len(colores_lugares)]
        folium.Marker(
            location=[lugar["lat"], lugar["lon"]], 
            popup=folium.Popup(f"<b>📍 {lugar['nombre']}</b>", max_width=300), 
            icon=folium.Icon(color=color, icon="star")
        ).add_to(mapa)
        folium.Circle(
            location=[lugar["lat"], lugar["lon"]], 
            radius=st.session_state.radio_km * 1000, 
            color=color, 
            fill=True, 
            fill_opacity=0.1,
            popup=f"Radio de búsqueda: {st.session_state.radio_km} km"
        ).add_to(mapa)
    
    # Configurar MarkerCluster con parámetros optimizados
    cluster = MarkerCluster(
        name="Espectaculares",
        options={
            'maxClusterRadius': 50,  # Radio máximo para clustering
            'disableClusteringAtZoom': 18,  # Desactivar clustering en zoom alto
            'showCoverageOnHover': True,
            'zoomToBoundsOnClick': True
        }
    ).add_to(mapa)
    
    # Agregar marcadores individualmente con mejor manejo de errores
    for i, r in df_filtrado.iterrows():
        try:
            # Validación exhaustiva de coordenadas
            lat = r["LATITUD"]
            lon = r["LONGITUD"]
            
            if (pd.isna(lat) or pd.isna(lon) or 
                not isinstance(lat, (int, float)) or 
                not isinstance(lon, (int, float)) or
                not (-90 <= lat <= 90) or 
                not (-180 <= lon <= 180)):
                marcadores_fallados += 1
                continue
            
            lugar_busqueda = r.get("LUGAR_BUSQUEDA", "Principal")
            color_index = next((idx for idx, lugar in enumerate(st.session_state.lugares_multiples) 
                              if lugar["nombre"] == lugar_busqueda), 0)
            color_marker = colores_lugares[color_index % len(colores_lugares)]
            
            # Crear contenido del popup más informativo
            popup_html = f"""
            <div style="min-width: 250px;">
                <h4 style="margin: 0; color: #333;">{r['CLAVE']}</h4>
                <hr style="margin: 5px 0;">
                <p style="margin: 2px 0;"><b>Distancia:</b> {r['DISTANCIA_KM']:.2f} km</p>
                <p style="margin: 2px 0;"><b>Tarifa:</b> {r['TARIFA_PUBLICO']}</p>
                <p style="margin: 2px 0;"><b>Lugar:</b> {lugar_busqueda}</p>
                <p style="margin: 2px 0;"><b>Tipo:</b> {r.get('TIPO', 'N/A')}</p>
                <p style="margin: 2px 0;"><b>Dirección:</b> {r.get('DIRECCION', 'N/A')}</p>
                <div style="margin-top: 8px;">
                    <a href='{r['MAPS_']}' target='_blank' style='color: blue; text-decoration: none;'>📍 Google Maps</a><br>
                    <a href='{r['STREET_VIEW']}' target='_blank' style='color: green; text-decoration: none;'>🌐 Street View</a>
                </div>
            </div>
            """
            
            folium.Marker(
                location=[lat, lon],
                popup=folium.Popup(popup_html, max_width=300),
                icon=folium.Icon(color=color_marker, icon="info-sign"),
                tooltip=f"{r['CLAVE']} - {r['TARIFA_PUBLICO']}"
            ).add_to(cluster)
            
            marcadores_agregados += 1
            
        except Exception as e:
            marcadores_fallados += 1
            # Mostrar errores específicos en modo debug
            # st.warning(f"Error en marcador {i}: {e}")
    
    # Agregar control de capas
    folium.LayerControl().add_to(mapa)
    
    # Mostrar estadísticas de marcadores
    st.info(f"**Estadísticas del mapa:** {marcadores_agregados} marcadores mostrados de {len(df_filtrado)} resultados totales")
    if marcadores_fallados > 0:
        st.warning(f"⚠️ {marcadores_fallados} marcadores no se pudieron mostrar debido a coordenadas inválidas o errores.")
    
    st.components.v1.html(folium.Figure().add_child(mapa).render(), height=500)

    # 5. SELECCIÓN DE ESPECTACULARES - SOLUCIÓN DEFINITIVA
    st.write("---")
    st.subheader("🎯 Selección de Espectaculares por Lugar")
    st.write("Selecciona los espectaculares que deseas incluir en las descargas.")
    
    # Función para manejar cambios en multiselect
    def actualizar_seleccion(lugar_nombre, nueva_seleccion):
        st.session_state.selecciones_por_lugar[lugar_nombre] = nueva_seleccion
        st.session_state.multiselect_actualizado = True
    
    # VERIFICACIÓN CORREGIDA: Usar st.session_state.df_por_lugar
    if st.session_state.df_por_lugar and len(st.session_state.df_por_lugar) > 0:
        tabs = st.tabs([f"📍 {lugar}" for lugar in st.session_state.df_por_lugar.keys()])
        
        todas_selecciones = []
        df_seleccionados_combinado = pd.DataFrame()
        
        for i, (lugar_nombre, df_lugar) in enumerate(st.session_state.df_por_lugar.items()):
            with tabs[i]:
                st.subheader(f"Espectaculares cerca de: {lugar_nombre}")
                
                if df_lugar.empty:
                    st.info(f"ℹ️ No se encontraron espectaculares cerca de {lugar_nombre}")
                    continue
                
                # Inicializar selección para este lugar si no existe
                if lugar_nombre not in st.session_state.selecciones_por_lugar:
                    st.session_state.selecciones_por_lugar[lugar_nombre] = []
                
                # Crear opciones con índices únicos
                opciones_lugar = [
                    f"{idx}. {row['CLAVE']} - {row['TARIFA_PUBLICO']} - {row['DISTANCIA_KM']} km - {row['TIPO']}" 
                    for idx, (_, row) in enumerate(df_lugar.iterrows())
                ]
                
                # Obtener selección actual
                seleccion_actual = st.session_state.selecciones_por_lugar.get(lugar_nombre, [])
                
                # Multiselect con callback
                seleccion_lugar = st.multiselect(
                    f"Selecciona espectaculares para {lugar_nombre}:",
                    options=opciones_lugar,
                    default=seleccion_actual,
                    key=f"multiselect_{lugar_nombre}_{i}"  # Clave única por lugar y pestaña
                )
                
                # ACTUALIZACIÓN INMEDIATA - SIN ESPERAR AL RERUN
                if seleccion_lugar != seleccion_actual:
                    actualizar_seleccion(lugar_nombre, seleccion_lugar)
                    st.rerun()
                
                # Mostrar resumen de selección
                if seleccion_lugar:
                    st.success(f"✅ **{len(seleccion_lugar)}** espectaculares seleccionados en {lugar_nombre}")
                    
                    # Procesar selección actual
                    indices_seleccionados_lugar = [int(op.split(".")[0]) for op in seleccion_lugar]
                    df_seleccionados_lugar = df_lugar.iloc[indices_seleccionados_lugar]
                    
                    # Mostrar tabla
                    columnas_mostrar = ['CLAVE', 'DIRECCION', 'TARIFA_PUBLICO', 'DISTANCIA_KM', 'TIPO']
                    st.dataframe(df_seleccionados_lugar[columnas_mostrar])
                    
                    # Agregar a combinación
                    todas_selecciones.extend(seleccion_lugar)
                    df_seleccionados_combinado = pd.concat([df_seleccionados_combinado, df_seleccionados_lugar])
                else:
                    st.info(f"ℹ️ No hay espectaculares seleccionados para {lugar_nombre}")
        
        # ACTUALIZAR SELECCIÓN GLOBAL
        st.session_state.espectaculares_seleccionados = todas_selecciones
        
        # 6. OPCIONES DE DESCARGA
        if todas_selecciones:
            st.write("---")
            st.subheader("💾 Opciones de Descarga")
            
            st.info(f"📋 **Folio de esta propuesta:** `{st.session_state.folio_actual}`")
            st.success(f"🎉 **Selección total:** {len(todas_selecciones)} espectaculares seleccionados de {len(st.session_state.df_por_lugar)} lugares")
            
            col_dl1, col_dl2 = st.columns(2)
            
            with col_dl1:
                csv_file = df_seleccionados_combinado.to_csv(index=False).encode('utf-8')
                if st.download_button(
                    label="⬇️ Descargar Resultados (CSV)", 
                    data=csv_file, 
                    file_name=f"{st.session_state.folio_actual}_resultados.csv", 
                    mime="text/csv",
                    key='download_csv'
                ):
                    incrementar_folio(st.session_state.lugares_multiples)
                    st.success(f"✅ Descarga completada.")
            
            with col_dl2:
                output = io.BytesIO()
                wb = Workbook()
                ws = wb.active
                ws.title = "Lugares cercanos"
                ws.append(list(df_seleccionados_combinado.columns))
                for r in df_seleccionados_combinado.to_dict('records'):
                    ws.append(list(r.values()))
                for row in range(2, ws.max_row + 1):
                    for col in range(1, ws.max_column + 1):
                        if ws.cell(row=1, column=col).value == "TARIFA_PUBLICO":
                            ws.cell(row=row, column=col).number_format = '"$"#,##0.00'
                            break
                tabla = Table(displayName="TablaEspectaculares", ref=f"A1:{get_column_letter(ws.max_column)}{ws.max_row}")
                tabla.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True)
                ws.add_table(tabla)
                for cell in ws[1]:
                    cell.fill = PatternFill(start_color="FFFF99", end_color="FFFF99", fill_type="solid")
                    cell.font = Font(bold=True)
                wb.save(output)
                if st.download_button(
                    label="⬇️ Descargar Resultados (Excel)", 
                    data=output.getvalue(), 
                    file_name=f"{st.session_state.folio_actual}_resultados.xlsx", 
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key='download_excel'
                ):
                    incrementar_folio(st.session_state.lugares_multiples)
                    st.success(f"✅ Descarga completada.")
            
            # 7. GENERAR PRESENTACIÓN - CORREGIDO
            st.write("---")
            st.subheader("🎓 Generar Presentación")
            
            if len(st.session_state.consultas_previas) > 1:
                combinar_consultas = st.checkbox(
                    "🔄 **Combinar con resultados de consultas anteriores**",
                    help="Incluir espectaculares de búsquedas previas en la misma presentación"
                )
            else:
                combinar_consultas = False
            
            if st.button("Crear Presentación", key='crear_presentacion'):
                try:
                    plantilla_pptx = "plantilla2.pptx"
                    prs = Presentation(plantilla_pptx)
                    if len(prs.slides) < 2:
                        st.error("❌ La plantilla debe tener al menos 2 diapositivas: la de título (0) y la de contenido (1).")
                    else:
                        slide_base_index = 1
                        slide_base = prs.slides[slide_base_index]
                        
                        if combinar_consultas:
                            todos_espectaculares = df_seleccionados_combinado.copy()
                            for consulta in st.session_state.consultas_previas[:-1]:
                                if not consulta["df_filtrado"].empty:
                                    todos_espectaculares = pd.concat([todos_espectaculares, consulta["df_filtrado"]], ignore_index=True)
                            
                            todos_espectaculares = todos_espectaculares.drop_duplicates(subset=['CLAVE'])
                            st.info(f"📊 Presentación combinada con {len(todos_espectaculares)} espectaculares de {len(st.session_state.consultas_previas)} consultas")
                        else:
                            todos_espectaculares = df_seleccionados_combinado
                        
                        for _, fila in todos_espectaculares.iterrows():
                            nueva_slide = duplicar_slide(prs, slide_base)
                            reemplazar_texto_slide(nueva_slide, fila.to_dict(), todos_espectaculares, st.session_state.nombre_negocio)
                        
                        # CORRECCIÓN: empty es un atributo, no un método
                        if not todos_espectaculares.empty:
                            primera_fila = todos_espectaculares.iloc[0]
                            reemplazar_texto_slide(prs.slides[0], primera_fila.to_dict(), todos_espectaculares, st.session_state.nombre_negocio)
                        
                        pptx_output = io.BytesIO()
                        prs.save(pptx_output)
                        
                        st.success(f"✅ ¡Presentación creada con éxito! - Folio: `{st.session_state.folio_actual}`")
                        
                        if st.download_button(
                            label="⬇️ Descargar Presentación (PPTX)", 
                            data=pptx_output.getvalue(), 
                            file_name=f"{st.session_state.folio_actual}.pptx", 
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key='download_pptx'
                        ):
                            incrementar_folio(st.session_state.lugares_multiples)
                            st.success(f"✅ Descarga completada.")
                except FileNotFoundError:
                    st.error("❌ **Error:** No se encontró el archivo de plantilla `plantilla2.pptx`. Asegúrate de que está en la misma carpeta que tu `app.py`.")
                except Exception as e:
                    st.error(f"❌ **Error al crear la presentación:** {e}")
        else:
            st.warning("⚠️ Por favor, selecciona al menos un espectacular de alguno de los lugares para habilitar las opciones de descarga.")
    else:
        st.info("ℹ️ No hay resultados separados por lugar para mostrar.")

# 8. HISTORIAL DE CONSULTAS
if len(st.session_state.consultas_previas) > 0:
    st.write("---")
    st.subheader("📊 Historial de Consultas")
    
    for i, consulta in enumerate(st.session_state.consultas_previas):
        with st.expander(f"🔍 Consulta {i+1}: {consulta['nombre_negocio']} - {consulta['fecha']} - {consulta['resultados']} resultados"):
            st.write(f"**Negocio:** {consulta['nombre_negocio']}")
            st.write(f"**Fecha:** {consulta['fecha']}")
            st.write(f"**Resultados:** {consulta['resultados']} espectaculares")
            if 'tipo' in consulta:
                st.write(f"**Tipo de búsqueda:** {consulta['tipo']}")
            if 'lugares' in consulta:
                st.write(f"**Lugares:** {', '.join(consulta['lugares'])}")
            if 'tipos_seleccionados' in consulta and consulta['tipos_seleccionados']:
                if len(consulta['tipos_seleccionados']) <= 5:
                    st.write(f"**Tipos seleccionados:** {', '.join(consulta['tipos_seleccionados'])}")
                else:
                    st.write(f"**Tipos seleccionados:** {len(consulta['tipos_seleccionados'])} tipos")





